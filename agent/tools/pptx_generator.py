import os
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

MIDNIGHT = RGBColor(0x02, 0x1C, 0x2E)
DEEP     = RGBColor(0x06, 0x5A, 0x82)
TEAL     = RGBColor(0x1C, 0x72, 0x93)
CYAN     = RGBColor(0x00, 0xD4, 0xFF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xE8, 0xF4, 0xF8)
GREY     = RGBColor(0xA0, 0xC4, 0xD8)
DARK_CARD= RGBColor(0x03, 0x35, 0x52)

SW = Inches(13.33)
SH = Inches(7.5)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _txt(slide, text, x, y, w, h, size, bold=False, italic=False,
         color=None, align=PP_ALIGN.LEFT, wrap=True):
    if not text:
        return
    box = slide.shapes.add_textbox(x, y, w, h)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else WHITE
    return box


def _title_slide(prs, title, subtitle):
    slide = _blank(prs)
    _bg(slide, MIDNIGHT)

    # Full-width deep blue top band
    _rect(slide, 0, 0, SW, Inches(0.22), CYAN)

    # Left panel
    _rect(slide, 0, Inches(0.22), Inches(6.5), SH - Inches(0.22), DEEP)

    # Right panel
    _rect(slide, Inches(6.5), Inches(0.22), Inches(6.83), SH - Inches(0.22), TEAL)

    # Title — left panel
    _txt(slide, title,
         Inches(0.5), Inches(1.5), Inches(5.7), Inches(4.0),
         34, bold=True, color=WHITE)

    # Cyan divider
    _rect(slide, Inches(0.5), Inches(5.6), Inches(3.0), Emu(50000), CYAN)

    # Subtitle
    _txt(slide, subtitle,
         Inches(0.5), Inches(5.8), Inches(5.7), Inches(0.8),
         13, italic=True, color=GREY)

    # Right panel — decorative large text
    _txt(slide, "AI",
         Inches(7.5), Inches(1.5), Inches(4.0), Inches(3.0),
         120, bold=True, color=RGBColor(0x04, 0x45, 0x65))
    _txt(slide, "Generated offline",
         Inches(7.2), Inches(5.5), Inches(5.5), Inches(0.6),
         14, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def _two_col_slide(prs, heading, bullets):
    slide = _blank(prs)
    _bg(slide, MIDNIGHT)

    _rect(slide, 0, 0, SW, Inches(1.0), DEEP)
    _rect(slide, 0, 0, Inches(0.22), Inches(1.0), CYAN)
    _txt(slide, heading,
         Inches(0.45), Inches(0.12), Inches(12.5), Inches(0.76),
         24, bold=True, color=WHITE)

    mid = (len(bullets) + 1) // 2
    left_b  = bullets[:mid]
    right_b = bullets[mid:]

    col_w = Inches(5.9)
    col_h = Inches(5.8)
    col_y = Inches(1.15)

    # Left card
    _rect(slide, Inches(0.45), col_y, col_w, col_h, DEEP)
    y = col_y + Inches(0.2)
    row_h = col_h / max(len(left_b), 1) - Inches(0.05)
    for b in left_b:
        _rect(slide, Inches(0.65), y + Inches(0.1), Emu(110000), Emu(110000), CYAN)
        _txt(slide, b, Inches(1.0), y, col_w - Inches(0.65), row_h,
             13, color=OFFWHITE)
        y += row_h + Inches(0.05)

    # Right card
    rx = Inches(0.45) + col_w + Inches(0.35)
    _rect(slide, rx, col_y, col_w, col_h, DARK_CARD)
    y = col_y + Inches(0.2)
    row_h = col_h / max(len(right_b), 1) - Inches(0.05)
    for b in right_b:
        _rect(slide, rx + Inches(0.2), y + Inches(0.1), Emu(110000), Emu(110000), CYAN)
        _txt(slide, b, rx + Inches(0.55), y, col_w - Inches(0.65), row_h,
             13, color=OFFWHITE)
        y += row_h + Inches(0.05)


def _strip_slide(prs, heading, bullets):
    slide = _blank(prs)
    _bg(slide, MIDNIGHT)

    _rect(slide, 0, 0, Inches(0.22), SH, DEEP)
    _txt(slide, heading,
         Inches(0.45), Inches(0.25), Inches(12.5), Inches(0.8),
         26, bold=True, color=CYAN)
    _rect(slide, Inches(0.45), Inches(1.1), Inches(3.5), Emu(45000), CYAN)

    n = len(bullets)
    available_h = SH - Inches(1.35) - Inches(0.15)
    row_h = available_h / max(n, 1)
    # Cap row height so text isn't too spread out
    row_h = min(row_h, Inches(0.82))
    strip_colors = [DEEP, DARK_CARD]
    y = Inches(1.3)

    for i, bullet in enumerate(bullets):
        if y + row_h > SH - Inches(0.1):
            break
        _rect(slide, Inches(0.45), y, Inches(12.43), row_h, strip_colors[i % 2])
        _rect(slide, Inches(0.45), y, Inches(0.55), row_h, CYAN)
        _txt(slide, str(i + 1),
             Inches(0.45), y + Inches(0.1), Inches(0.55), row_h - Inches(0.1),
             16, bold=True, color=MIDNIGHT, align=PP_ALIGN.CENTER)
        _txt(slide, bullet,
             Inches(1.1), y + Inches(0.08),
             Inches(11.7), row_h - Inches(0.1),
             13, color=OFFWHITE)
        y += row_h + Inches(0.02)


def _card_grid_slide(prs, heading, bullets):
    slide = _blank(prs)
    _bg(slide, MIDNIGHT)

    _rect(slide, 0, 0, SW, Inches(1.0), DEEP)
    _rect(slide, 0, 0, Inches(0.22), Inches(1.0), CYAN)
    _txt(slide, heading,
         Inches(0.45), Inches(0.12), Inches(12.5), Inches(0.76),
         24, bold=True, color=WHITE)

    n = len(bullets)
    cols = 4 if n >= 6 else 3 if n >= 3 else n
    rows = (n + cols - 1) // cols

    pad = Inches(0.2)
    total_w = SW - Inches(0.9)
    total_h = SH - Inches(1.35)
    card_w = (total_w - pad * (cols - 1)) / cols
    card_h = (total_h - pad * (rows - 1)) / rows

    card_colors = [DEEP, TEAL, DARK_CARD]
    icons = ["▸", "◆", "●", "★", "▲", "■", "▶", "◉"]

    for i, bullet in enumerate(bullets):
        row = i // cols
        col = i % cols
        cx = Inches(0.45) + col * (card_w + pad)
        cy = Inches(1.2) + row * (card_h + pad)

        _rect(slide, cx, cy, card_w, card_h, card_colors[i % 3])

        circ = Inches(0.35)
        s = slide.shapes.add_shape(9, cx + Inches(0.12), cy + Inches(0.1), circ, circ)
        s.fill.solid()
        s.fill.fore_color.rgb = CYAN
        s.line.fill.background()

        _txt(slide, icons[i % len(icons)],
             cx + Inches(0.12), cy + Inches(0.1), circ, circ,
             11, bold=True, color=MIDNIGHT, align=PP_ALIGN.CENTER)

        _txt(slide, bullet,
             cx + Inches(0.12), cy + Inches(0.52),
             card_w - Inches(0.24), card_h - Inches(0.62),
             12, color=OFFWHITE, wrap=True)


def _closing_slide(prs, heading, bullets):
    slide = _blank(prs)
    _bg(slide, MIDNIGHT)

    _rect(slide, 0, 0, SW, Inches(1.8), DEEP)
    _txt(slide, heading,
         Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.2),
         30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(4.5), Inches(1.85), Inches(4.33), Emu(50000), CYAN)

    # All bullets as two columns
    mid = (len(bullets) + 1) // 2
    left_b  = bullets[:mid]
    right_b = bullets[mid:]

    col_w = Inches(5.9)
    col_y = Inches(2.05)
    col_h = SH - col_y - Inches(0.2)
    row_h = col_h / max(len(left_b), 1)
    row_h = min(row_h, Inches(0.75))

    # Left
    y = col_y
    for b in left_b:
        _rect(slide, Inches(0.45), y + Inches(0.05), Emu(110000), Emu(110000), CYAN)
        _txt(slide, b, Inches(0.85), y, col_w - Inches(0.5), row_h,
             13, color=OFFWHITE)
        y += row_h

    # Right
    rx = Inches(0.45) + col_w + Inches(0.5)
    y = col_y
    row_h_r = col_h / max(len(right_b), 1)
    row_h_r = min(row_h_r, Inches(0.75))
    for b in right_b:
        _rect(slide, rx, y + Inches(0.05), Emu(110000), Emu(110000), CYAN)
        _txt(slide, b, rx + Inches(0.4), y, col_w - Inches(0.5), row_h_r,
             13, color=OFFWHITE)
        y += row_h_r


# ── Public API ─────────────────────────────────────────────────────────────────

def generate_ppt(title: str, slides: list, output_path: str = None,
                 subtitle: str = "") -> dict:
    if not PPTX_AVAILABLE:
        return {"success": False, "path": "", "slide_count": 0,
                "message": "python-pptx not installed. Run: pip install python-pptx"}

    if output_path is None:
        safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:50]
        ts = datetime.now().strftime("%H%M%S")
        out_dir = os.path.abspath(
            os.path.join(os.path.dirname(__file__), "..", "..", "demo", "outputs"))
        os.makedirs(out_dir, exist_ok=True)
        output_path = os.path.join(out_dir, f"{safe}_{ts}.pptx")

    try:
        prs = Presentation()
        prs.slide_width  = SW
        prs.slide_height = SH

        _title_slide(prs, title, subtitle or "Generated by Offline AI OS")

        content_layouts = [_two_col_slide, _card_grid_slide, _strip_slide]

        for i, sd in enumerate(slides):
            h = sd.get("heading", f"Section {i+1}")
            b = sd.get("bullets", [])
            if not b:
                b = ["No content available."]

            if i == len(slides) - 1:
                _closing_slide(prs, h, b)
            else:
                content_layouts[i % len(content_layouts)](prs, h, b)

        prs.save(output_path)
        return {"success": True, "path": os.path.abspath(output_path),
                "slide_count": len(slides) + 1,
                "message": f"Saved: {output_path}"}

    except Exception as e:
        return {"success": False, "path": "", "slide_count": 0, "message": str(e)}
